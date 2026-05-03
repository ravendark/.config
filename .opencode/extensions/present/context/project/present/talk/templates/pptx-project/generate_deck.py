#!/usr/bin/env python3
"""
generate_deck.py -- Theme-aware PowerPoint deck generator.

This is a skeleton script demonstrating how to assemble a PPTX presentation
using python-pptx with the talk library's theme system. The pptx-assembly-agent
uses these patterns when producing PowerPoint output.

Usage:
    python generate_deck.py --theme academic-clean --output talk.pptx

Requirements:
    pip install python-pptx
"""

import argparse
import json
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Theme loading
# ---------------------------------------------------------------------------

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
THEME_FILE = os.path.join(SCRIPT_DIR, "theme_mappings.json")

VALID_THEMES = ["academic-clean", "clinical-teal", "ucsf-institutional"]


def load_themes():
    """Load theme mappings from theme_mappings.json."""
    with open(THEME_FILE) as f:
        return json.load(f)


def hex_to_rgb(hex_str):
    """Convert '#RRGGBB' hex string to pptx RGBColor."""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------

def add_blank_slide(prs):
    """Add a blank slide (layout index 6)."""
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


def add_titled_slide(prs, theme, title_text):
    """Add a blank slide with a themed title at the top."""
    slide = add_blank_slide(prs)
    typo = theme["typography"]
    palette = theme["palette"]

    txBox = slide.shapes.add_textbox(
        left=Inches(0.75), top=Inches(0.4),
        width=Inches(11.8), height=Inches(0.8),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.name = typo["heading_font"]
    run.font.size = Pt(typo["heading_size_pt"])
    run.font.bold = typo["heading_bold"]
    run.font.color.rgb = hex_to_rgb(palette["heading"])

    return slide


def apply_body_style(run, theme):
    """Apply body text styling from theme."""
    typo = theme["typography"]
    run.font.name = typo["body_font"]
    run.font.size = Pt(typo["body_size_pt"])
    run.font.bold = False
    run.font.color.rgb = hex_to_rgb(theme["palette"]["text"])


def apply_caption_style(run, theme):
    """Apply caption styling from theme."""
    typo = theme["typography"]
    run.font.name = typo["body_font"]
    run.font.size = Pt(typo["caption_size_pt"])
    run.font.color.rgb = hex_to_rgb(theme["palette"]["muted"])


def set_speaker_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


# ---------------------------------------------------------------------------
# Component helpers (see pptx-generation.md for full documentation)
# ---------------------------------------------------------------------------

def add_pptx_table(slide, theme, headers, rows, caption=None,
                   highlight_row=None):
    """Add a themed data table. See pptx-generation.md section 4.1."""
    palette = theme["palette"]
    typo = theme["typography"]
    sp = theme["spacing"]

    left = Inches(sp["content_left"])
    top = Inches(2.0)
    width = Inches(sp["content_width"])
    n_rows = len(rows) + 1
    n_cols = len(headers)
    row_h = Inches(0.45)
    height = row_h * n_rows

    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = shape.table

    for ci, header in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = hex_to_rgb(palette["accent_light"])
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.name = typo["body_font"]
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = hex_to_rgb(palette["heading"])

    for ri, row_data in enumerate(rows):
        for ci, value in enumerate(row_data):
            cell = table.cell(ri + 1, ci)
            cell.text = str(value)
            if highlight_row is not None and ri == highlight_row:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hex_to_rgb(palette["accent_light"])
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = typo["body_font"]
                    run.font.size = Pt(13)
                    run.font.color.rgb = hex_to_rgb(palette["text"])

    if caption:
        cap = slide.shapes.add_textbox(
            left, top + height + Inches(0.1), width, Inches(0.4)
        )
        p = cap.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = caption
        apply_caption_style(run, theme)

    return shape


def add_pptx_figure(slide, theme, image_path, caption=None, source=None):
    """Add a figure with caption. See pptx-generation.md section 4.2."""
    sp = theme["spacing"]
    left = Inches(sp["content_left"])
    top = Inches(1.8)
    max_w = Inches(10.0)

    pic = slide.shapes.add_picture(image_path, left, top, width=max_w)

    img_bottom = pic.top + pic.height
    if caption:
        cap = slide.shapes.add_textbox(
            left, img_bottom + Inches(0.15),
            Inches(sp["content_width"]), Inches(0.4),
        )
        p = cap.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = caption
        apply_caption_style(run, theme)
        img_bottom = cap.top + cap.height

    if source:
        src = slide.shapes.add_textbox(
            left, img_bottom + Inches(0.05),
            Inches(sp["content_width"]), Inches(0.3),
        )
        p = src.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"Source: {source}"
        run.font.name = theme["typography"]["body_font"]
        run.font.size = Pt(11)
        run.font.italic = True
        run.font.color.rgb = hex_to_rgb(theme["palette"]["muted"])

    return pic


# ---------------------------------------------------------------------------
# Example deck
# ---------------------------------------------------------------------------

def build_example_deck(theme_name, output_path):
    """Build an example presentation demonstrating all slide types."""
    themes = load_themes()
    if theme_name not in themes:
        print(f"Error: Unknown theme '{theme_name}'. Choose from: {VALID_THEMES}")
        sys.exit(1)

    theme = themes[theme_name]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Title slide ---
    slide = add_titled_slide(prs, theme, "Example Research Presentation")
    txBox = slide.shapes.add_textbox(
        Inches(0.75), Inches(1.5), Inches(11.8), Inches(1.0)
    )
    p = txBox.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Author Name | Institution | Date"
    apply_body_style(run, theme)
    set_speaker_notes(slide, "Welcome and introduction.\n- Introduce yourself\n- State the research question")

    # --- Content slide with bullets ---
    slide = add_titled_slide(prs, theme, "Background")
    bullets = [
        "Key finding from prior literature",
        "Gap in current knowledge",
        "Our approach to addressing this gap",
    ]
    for i, bullet in enumerate(bullets):
        txBox = slide.shapes.add_textbox(
            Inches(1.0), Inches(1.6 + i * 0.5),
            Inches(11.0), Inches(0.45),
        )
        p = txBox.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"\u2022  {bullet}"
        apply_body_style(run, theme)
    set_speaker_notes(slide, "Provide context for the study.\n- Cite 2-3 key papers")

    # --- Table slide ---
    slide = add_titled_slide(prs, theme, "Results: Baseline Characteristics")
    add_pptx_table(
        slide, theme,
        headers=["Variable", "Group A (n=100)", "Group B (n=100)", "p-value"],
        rows=[
            ["Age, mean (SD)", "55.2 (12.1)", "54.8 (11.9)", "0.82"],
            ["Female, n (%)", "48 (48%)", "52 (52%)", "0.57"],
            ["BMI, mean (SD)", "27.3 (4.2)", "27.1 (4.0)", "0.73"],
            ["Diabetes, n (%)", "22 (22%)", "25 (25%)", "0.63"],
        ],
        caption="Table 1. Baseline characteristics by treatment group.",
    )
    set_speaker_notes(slide, "Groups are well-balanced at baseline.\n- No significant differences")

    # --- Save ---
    prs.save(output_path)
    print(f"Saved: {output_path} ({len(prs.slides)} slides, theme: {theme_name})")


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="Generate themed PPTX deck")
    parser.add_argument(
        "--theme", default="academic-clean", choices=VALID_THEMES,
        help="Theme name (default: academic-clean)",
    )
    parser.add_argument(
        "--output", default="output.pptx",
        help="Output file path (default: output.pptx)",
    )
    args = parser.parse_args()
    build_example_deck(args.theme, args.output)


if __name__ == "__main__":
    main()
